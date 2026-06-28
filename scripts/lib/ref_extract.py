"""Extract outbound links and embedded images from fetched Office files.

The depth guarantee of the crawl lives here: every link a file contains must
be surfaced so the frontier loop can enqueue it. For spreadsheets that means
every tab and every cell hyperlink; for decks, every slide and shape.

Return contract (R3):
    extract_links_and_images(local_path, image_dir)
        -> tuple[list[str], list[dict]]

    urls:   list of external URL strings.
    images: list of dicts, one per embedded image written, in document order:
            {"path": <str>, "section": <str|None>, "nearby_text": <str|None>}

    - section    : text of the nearest *preceding* heading paragraph
                   (style name starts with "Heading"); None if not found.
    - nearby_text: text of the nearest *preceding* non-empty non-heading
                   paragraph; None if not found.

Tasks 4/5 append .xlsx/.pptx branches and map their own context onto the
same three keys — do NOT rename them.
"""
from __future__ import annotations

import pathlib
import re as _re

_HYPERLINK_REL = "hyperlink"
_IMAGE_REL = "image"

# Matches the URL inside =HYPERLINK("url", ...) formula strings.
_HYPERLINK_FORMULA = _re.compile(r'HYPERLINK\(\s*"([^"]+)"', _re.IGNORECASE)
# Matches bare https?:// URLs pasted as visible text in cells.
_BARE_URL = _re.compile(r'https?://[^\s"\'<>)\]]+')


def _write_blob(blob: bytes, image_dir: str, stem: str, idx: int, ext: str) -> str:
    d = pathlib.Path(image_dir)
    d.mkdir(parents=True, exist_ok=True)
    out = d / f"{stem}__{idx}{ext}"
    out.write_bytes(blob)
    return str(out)


def _extract_docx(path: str, image_dir: str) -> tuple[list[str], list[dict]]:
    """Extract URLs and context-tagged images from a .docx file.

    URL extraction: walks doc.part.rels for external hyperlink relationships
    (these have no document position but are exhaustive for all hyperlinks).

    Image extraction: walks the document body in element order, tracking the
    current heading and last non-empty paragraph as context.  Inline images
    are found via a:blip/@r:embed → the image part's blob.  This body-walk
    technique is adapted from raw/se-runbook/_extract/extract_maindoc.py.txt.
    """
    from docx import Document
    from docx.oxml.ns import qn
    from docx.text.paragraph import Paragraph

    doc = Document(path)
    stem = pathlib.Path(path).stem

    # ── 1. Collect URLs from all relationships ────────────────────────────
    urls: list[str] = []
    for rel in doc.part.rels.values():
        if _HYPERLINK_REL in rel.reltype and rel.is_external:
            urls.append(rel.target_ref)

    # ── 2. Body-walk for images with section/nearby_text context ─────────
    images: list[dict] = []
    idx = 0
    current_section: str | None = None
    last_nearby: str | None = None

    for child in doc.element.body.iterchildren():
        if child.tag == qn("w:p"):
            p = Paragraph(child, doc)
            style_name: str = p.style.name if p.style is not None else ""
            text = p.text.strip()

            if style_name.startswith("Heading"):
                current_section = text or None
                # A heading also resets nearby_text — the heading IS the new
                # context boundary; paragraph text after the heading is what
                # we want for nearby_text, not the heading itself.
                last_nearby = None
            elif text:
                last_nearby = text

            # Find all inline image references in this paragraph
            for blip in p._p.findall(".//" + qn("a:blip")):
                embed = blip.get(qn("r:embed"))
                if embed and embed in doc.part.rels:
                    rel = doc.part.rels[embed]
                    try:
                        blob = rel.target_part.blob
                        ext = pathlib.Path(rel.target_part.partname).suffix or ".png"
                        img_path = _write_blob(blob, image_dir, stem, idx, ext)
                        images.append(
                            {
                                "path": img_path,
                                "section": current_section,
                                "nearby_text": last_nearby,
                            }
                        )
                        idx += 1
                    except Exception:
                        pass

    return urls, images


def _extract_xlsx(path: str, image_dir: str) -> tuple[list[str], list[dict]]:
    """Extract URLs and images from a .xlsx file.

    URL extraction: walks EVERY worksheet and collects three forms of link:
      1. cell.hyperlink.target  — structured hyperlink relationship on the cell
      2. =HYPERLINK("url",...)  — formula string (Google Sheets export common case)
      3. bare https?:// URL     — pasted-as-text in a cell value (most common real case)

    Image extraction: reads ws._images on each worksheet; each image is returned
    as the locked dict shape {"path", "section", "nearby_text"} where section is
    the worksheet tab name (R3 mapping: xlsx tab_name → section) and nearby_text
    is None (robustly reading anchor-adjacent cells is fiddly; None is acceptable
    per spec and consistent with YAGNI).
    """
    import openpyxl

    wb = openpyxl.load_workbook(path)  # data_only=False (default) — preserves formulas
    stem = pathlib.Path(path).stem
    urls: list[str] = []
    images: list[dict] = []
    idx = 0

    for ws in wb.worksheets:            # EVERY tab
        for row in ws.iter_rows():
            for cell in row:
                # Form 1: structured hyperlink on the cell
                if cell.hyperlink and cell.hyperlink.target:
                    urls.append(cell.hyperlink.target)
                # Form 2 & 3: scan string cell values
                if isinstance(cell.value, str):
                    if "HYPERLINK(" in cell.value.upper():
                        urls.extend(_HYPERLINK_FORMULA.findall(cell.value))
                    urls.extend(_BARE_URL.findall(cell.value))

        # Images: ws._images is the list openpyxl populates from drawings
        for img in getattr(ws, "_images", []):
            try:
                blob = img._data()
                images.append({
                    "path": _write_blob(blob, image_dir, stem, idx, ".png"),
                    "section": ws.title,   # R3: xlsx tab_name → section
                    "nearby_text": None,   # acceptable per spec (YAGNI)
                })
                idx += 1
            except Exception:
                pass

    return urls, images


def _extract_pptx(path: str, image_dir: str) -> tuple[list[str], list[dict]]:
    """Extract URLs and context-tagged images from a .pptx file.

    URL extraction: for every slide, recurses into GROUP shapes and collects:
      1. shape click_action.hyperlink.address  — shape-level click hyperlink
      2. text-run run.hyperlink.address        — hyperlink on an individual run
      3. bare https?:// URLs in shape text     — pasted-as-text (reuses _BARE_URL)

    Image extraction: every PICTURE shape's blob is written to image_dir.
    Context (R3 mapping):
      - section     = slide title placeholder text (slide_title), or None
      - nearby_text = all non-empty text on the slide joined by newline, or None
    Both values are computed ONCE per slide (including recursing groups) and
    tagged onto every picture found on that slide, at any nesting depth.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(path)
    stem = pathlib.Path(path).stem
    urls: list[str] = []
    images: list[dict] = []
    counter = [0]  # mutable cell so nested walker can bump it

    def _gather_text(shapes) -> list[str]:
        """Recursively collect all non-empty text-frame texts across shapes."""
        texts: list[str] = []
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                texts.extend(_gather_text(shape.shapes))
            elif shape.has_text_frame:
                t = (shape.text_frame.text or "").strip()
                if t:
                    texts.append(t)
        return texts

    def walk(shapes, slide_title: str | None, slide_text: str | None) -> None:
        for shape in shapes:
            # Shape-level click hyperlink
            try:
                addr = shape.click_action.hyperlink.address
                if addr:
                    urls.append(addr)
            except Exception:
                pass

            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                # Recurse — pass slide-level context unchanged
                walk(shape.shapes, slide_title, slide_text)
                continue

            # Text-run hyperlinks + bare URLs
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            addr = run.hyperlink.address
                            if addr:
                                urls.append(addr)
                        except Exception:
                            pass
                urls.extend(_BARE_URL.findall(shape.text_frame.text or ""))

            # Pictures — write blob and return R3 dict with slide context
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    ext = "." + (image.ext or "png")
                    images.append({
                        "path": _write_blob(image.blob, image_dir, stem, counter[0], ext),
                        "section": slide_title,    # R3: pptx slide_title → section
                        "nearby_text": slide_text, # R3: pptx slide_text  → nearby_text
                    })
                    counter[0] += 1
                except Exception:
                    pass

    for slide in prs.slides:
        # Compute slide-level context once (including text in groups)
        if slide.shapes.title is not None:
            slide_title: str | None = (slide.shapes.title.text or "").strip() or None
        else:
            slide_title = None

        gathered = _gather_text(slide.shapes)
        slide_text: str | None = "\n".join(gathered) or None

        walk(slide.shapes, slide_title, slide_text)

    return urls, images


def extract_links_and_images(
    local_path: str, image_dir: str
) -> tuple[list[str], list[dict]]:
    """Dispatch to the correct extractor based on file extension.

    Supported: .docx (Task 3).  .xlsx and .pptx added in Tasks 4/5.

    Raises:
        ValueError: if the file extension has no registered extractor.
    """
    ext = pathlib.Path(local_path).suffix.lower()
    if ext == ".docx":
        return _extract_docx(local_path, image_dir)
    if ext == ".xlsx":
        return _extract_xlsx(local_path, image_dir)
    if ext == ".pptx":
        return _extract_pptx(local_path, image_dir)
    raise ValueError(f"no extractor for {ext!r}")


def extract_pdf_links(pdf_path: str) -> list:
    """Scrape annotation URIs + visible-text URLs from a PDF render. Never raises."""
    if not pathlib.Path(pdf_path).exists():
        return []
    urls = []
    try:
        import pdfplumber

        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                for annot in (page.annots or []):
                    uri = annot.get("uri") or (annot.get("data") or {}).get("A", {}).get("URI")
                    if uri:
                        urls.append(uri)
                urls.extend(_BARE_URL.findall(page.extract_text() or ""))
    except Exception:
        return urls
    return urls
