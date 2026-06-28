"""Tests for scripts/lib/ref_extract.py — Task 3 (docx) with R3 context."""
from __future__ import annotations

import struct
import zlib
from io import BytesIO

import pytest
from docx import Document

from scripts.lib.ref_extract import extract_links_and_images


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _minimal_png() -> bytes:
    """Return the bytes of a minimal valid 1×1 white PNG."""
    def _chunk(name: bytes, data: bytes) -> bytes:
        c = zlib.crc32(name + data) & 0xFFFFFFFF
        return struct.pack(">I", len(data)) + name + data + struct.pack(">I", c)

    sig = b"\x89PNG\r\n\x1a\n"
    ihdr_data = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
    ihdr = _chunk(b"IHDR", ihdr_data)
    raw_row = b"\x00\xFF\xFF\xFF"          # filter byte + 3 RGB bytes
    idat = _chunk(b"IDAT", zlib.compress(raw_row))
    iend = _chunk(b"IEND", b"")
    return sig + ihdr + idat + iend


# ---------------------------------------------------------------------------
# Test 1 — hyperlink relationship surfaces in urls
# ---------------------------------------------------------------------------

def test_docx_links_extracted(tmp_path):
    doc = Document()
    p = doc.add_paragraph()
    p.add_run("see this")
    part = doc.part
    part.relate_to(
        "https://docs.google.com/spreadsheets/d/SHEET1/edit",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
        is_external=True,
    )
    path = tmp_path / "doc.docx"
    doc.save(str(path))

    urls, images = extract_links_and_images(str(path), str(tmp_path / "img"))
    assert "https://docs.google.com/spreadsheets/d/SHEET1/edit" in urls


# ---------------------------------------------------------------------------
# Test 2 — R3 context: section + nearby_text on returned image dicts
# ---------------------------------------------------------------------------

def test_docx_image_context(tmp_path):
    """Image dict must carry section (nearest preceding heading) and
    nearby_text (nearest preceding non-empty paragraph)."""
    doc = Document()
    doc.add_heading("Setup", level=1)
    doc.add_paragraph("Open the admin console")
    # Add an inline picture — python-docx writes it as a Run with w:drawing
    doc.add_picture(BytesIO(_minimal_png()))

    path = tmp_path / "ctx.docx"
    doc.save(str(path))

    img_dir = tmp_path / "img"
    urls, images = extract_links_and_images(str(path), str(img_dir))

    assert len(images) == 1
    img = images[0]
    assert img["section"] == "Setup"
    assert "Open the admin console" in img["nearby_text"]
    # The written file must actually exist
    import pathlib
    assert pathlib.Path(img["path"]).exists()


# ---------------------------------------------------------------------------
# Test 3 — dict has required keys (path / section / nearby_text)
# ---------------------------------------------------------------------------

def test_docx_image_dict_shape(tmp_path):
    doc = Document()
    doc.add_picture(BytesIO(_minimal_png()))
    path = tmp_path / "shape.docx"
    doc.save(str(path))

    _, images = extract_links_and_images(str(path), str(tmp_path / "img"))
    assert len(images) == 1
    img = images[0]
    assert set(img.keys()) == {"path", "section", "nearby_text"}


# ---------------------------------------------------------------------------
# Test 4 — None context when no heading / paragraph precede the image
# ---------------------------------------------------------------------------

def test_docx_image_no_context(tmp_path):
    """When no heading or paragraph precede the image, both fields are None."""
    doc = Document()
    doc.add_picture(BytesIO(_minimal_png()))
    path = tmp_path / "nocontext.docx"
    doc.save(str(path))

    _, images = extract_links_and_images(str(path), str(tmp_path / "img"))
    assert len(images) == 1
    img = images[0]
    assert img["section"] is None
    assert img["nearby_text"] is None


# ---------------------------------------------------------------------------
# Test 5 — dispatcher raises ValueError for unsupported extension
# ---------------------------------------------------------------------------

def test_unsupported_extension_raises(tmp_path):
    fake = tmp_path / "note.txt"
    fake.write_text("hello")
    with pytest.raises(ValueError, match="no extractor for"):
        extract_links_and_images(str(fake), str(tmp_path / "img"))


# ---------------------------------------------------------------------------
# Test 6 — heading immediately before image (no paragraph in between)
# ---------------------------------------------------------------------------

def test_docx_image_directly_after_heading_no_nearby(tmp_path):
    """When a heading directly precedes an image with no paragraph in between,
    nearby_text is None (heading is the boundary; text lives in section)."""
    doc = Document()
    doc.add_heading("Config", level=2)
    doc.add_picture(BytesIO(_minimal_png()))
    path = tmp_path / "headingonly.docx"
    doc.save(str(path))

    _, images = extract_links_and_images(str(path), str(tmp_path / "img"))
    assert len(images) == 1
    assert images[0]["section"] == "Config"
    assert images[0]["nearby_text"] is None


# ---------------------------------------------------------------------------
# XLSX Tests (Task 4)
# ---------------------------------------------------------------------------

import openpyxl


def test_xlsx_walks_all_tabs_and_formula_links(tmp_path):
    wb = openpyxl.Workbook()
    ws1 = wb.active
    ws1.title = "Tab1"
    ws1["A1"].value = "linked"
    ws1["A1"].hyperlink = "https://docs.google.com/document/d/DOC_IN_TAB1/edit"
    ws2 = wb.create_sheet("Tab2")
    ws2["B2"].value = '=HYPERLINK("https://docs.google.com/presentation/d/DECK_IN_TAB2/edit","go")'
    ws2["B3"].value = "see https://docs.google.com/document/d/PASTED_AS_TEXT/edit for details"
    path = tmp_path / "book.xlsx"
    wb.save(str(path))

    urls, _ = extract_links_and_images(str(path), str(tmp_path / "img"))
    assert "https://docs.google.com/document/d/DOC_IN_TAB1/edit" in urls          # cell hyperlink, tab 1
    assert "https://docs.google.com/presentation/d/DECK_IN_TAB2/edit" in urls     # =HYPERLINK formula, tab 2
    assert "https://docs.google.com/document/d/PASTED_AS_TEXT/edit" in urls       # bare URL in text, tab 2


def test_xlsx_image_dict_shape_and_section(tmp_path):
    """Image from xlsx must return the locked dict shape with section=tab_name."""
    import pathlib
    from openpyxl.drawing.image import Image as XLImage

    # Write the minimal PNG to a temp file (openpyxl.drawing.image.Image needs a path/file)
    png_path = tmp_path / "tiny.png"
    png_path.write_bytes(_minimal_png())

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Diagrams"
    img = XLImage(str(png_path))
    ws.add_image(img, "A1")
    path = tmp_path / "imgbook.xlsx"
    wb.save(str(path))

    img_dir = tmp_path / "img"
    _, images = extract_links_and_images(str(path), str(img_dir))

    assert len(images) == 1
    result = images[0]
    assert set(result.keys()) == {"path", "section", "nearby_text"}
    assert result["section"] == "Diagrams"
    assert result["nearby_text"] is None
    assert pathlib.Path(result["path"]).exists()


# ---------------------------------------------------------------------------
# PPTX Tests (Task 5)
# ---------------------------------------------------------------------------

from pptx import Presentation
from pptx.util import Inches


def test_pptx_extracts_run_hyperlink(tmp_path):
    """Run hyperlink inside a text-run must surface in urls list."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = "open sheet"
    run.hyperlink.address = "https://docs.google.com/spreadsheets/d/SHEET_IN_DECK/edit"
    path = tmp_path / "deck.pptx"
    prs.save(str(path))

    urls, _ = extract_links_and_images(str(path), str(tmp_path / "img"))
    assert "https://docs.google.com/spreadsheets/d/SHEET_IN_DECK/edit" in urls


def test_pptx_picture_dict_shape_and_slide_context(tmp_path):
    """Picture from pptx must return the locked dict with section=slide_title,
    nearby_text containing body text on that slide."""
    import pathlib

    prs = Presentation()
    # slide_layouts[5] = "Title Only" — provides a title placeholder
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Setup Slide"

    # Add a textbox with body text
    tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    tb.text_frame.paragraphs[0].add_run().text = "Click the gear icon"

    # Add a picture
    slide.shapes.add_picture(BytesIO(_minimal_png()), Inches(1), Inches(3))

    path = tmp_path / "ctx_deck.pptx"
    prs.save(str(path))

    img_dir = tmp_path / "img"
    _, images = extract_links_and_images(str(path), str(img_dir))

    assert len(images) == 1
    img = images[0]
    assert set(img.keys()) == {"path", "section", "nearby_text"}
    assert img["section"] == "Setup Slide"
    assert "Click the gear icon" in (img["nearby_text"] or "")
    assert pathlib.Path(img["path"]).exists()


# ---------------------------------------------------------------------------
# PDF Tests (Task 7b)
# ---------------------------------------------------------------------------

import os
import pytest
from scripts.lib.ref_extract import extract_pdf_links


def test_extract_pdf_links_returns_list(tmp_path):
    # Build a tiny PDF with a visible URL using pdfplumber's dependency (pdfminer can't write);
    # instead assert the function handles a missing/empty file gracefully and returns a list.
    missing = str(tmp_path / "nope.pdf")
    assert extract_pdf_links(missing) == []
